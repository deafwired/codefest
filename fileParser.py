import os
import sys
import pymupdf
from docx import Document
from pptx import Presentation
from util.fileTypes import FileTypeClassifier

def classify(self, file_path):
    ext = os.path.splitext(file_path)[1].lower()
    return self.file_types.get(ext, "unknown")

classifier = FileTypeClassifier()

def extract_text_from_docx(file_path):
    doc = Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_xlsx(file_path):
    df = pd.read_excel(file_path, engine="openpyxl")
    return df.to_string(index=False)

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_text_from_pdf(file_path):
    doc = pymupdf.open(file_path)
    text = "\n".join([page.get_text() for page in doc])
    return text if text else "No extractable text found."

def extract_text_default(file_path):
    extracted_text = ""
    try:
        with open(file_path, "r") as file:
            if file is None:
                return None
            for line in file:
                extracted_text += line.strip("\n") + " "
        return extracted_text.strip()
    except FileNotFoundError:
        print(f"Error: The file '{file_path}' was not found.")
        return None
    except IOError:
        print("Unsupported file type.")
        return None




def extract_text(file_path):
    file_extension = os.path.splitext(file_path)[1]
    print(file_extension)

    if file_extension == ".docx":
        return extract_text_from_docx(file_path)
    elif file_extension == ".xlsx":
        return extract_text_from_xlsx(file_path)
    elif file_extension == ".pptx":
        return extract_text_from_pptx(file_path)
    elif file_extension == ".pdf":
        return extract_text_from_pdf(file_path)
    else:
        return extract_text_default(file_path)

if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("Usage: python script.py <file_path>")
        sys.exit(1)

    file_path = sys.argv[1]
    if not os.path.exists(file_path):
        print("File not found.")
        sys.exit(1)

    extracted_text = extract_text(file_path)
    print("\nExtracted Content:\n")
    print(extracted_text)